"""
File processing service for architecture documents
"""

import os
import logging
import asyncio
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import uuid

# File processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF for PDF processing
import pytesseract
from io import BytesIO

from app.models.architecture import (
    ArchitectureFile, FileType, FileUploadStatus, FileProcessingResult,
    ArchitectureDomain, ArchitectureComponent, ArchitecturePattern, ArchitectureDecision
)
from app.config import settings

logger = logging.getLogger(__name__)


class FileProcessor:
    """Service for processing architecture files"""
    
    def __init__(self):
        self.supported_types = {
            FileType.PDF: self._process_pdf,
            FileType.DOCX: self._process_docx,
            FileType.PPTX: self._process_pptx,
            FileType.PNG: self._process_image,
            FileType.JPG: self._process_image,
            FileType.JPEG: self._process_image,
            FileType.SVG: self._process_svg,
        }
        
        # Create upload directory if it doesn't exist
        self.upload_dir = Path("uploads")
        self.upload_dir.mkdir(exist_ok=True)
        
        # Create processed files directory
        self.processed_dir = Path("processed")
        self.processed_dir.mkdir(exist_ok=True)
    
    async def process_file(self, file_path: str, file_type: FileType, file_id: str) -> FileProcessingResult:
        """Process an uploaded file"""
        start_time = datetime.utcnow()
        
        try:
            logger.info(f"Processing file {file_id} of type {file_type}")
            
            # Get the appropriate processor
            processor = self.supported_types.get(file_type)
            if not processor:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            # Process the file
            result = await processor(file_path, file_id)
            
            # Calculate processing time
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            
            # Create processing result
            processing_result = FileProcessingResult(
                file_id=file_id,
                processing_success=True,
                processing_errors=[],
                text_content=result.get("text_content"),
                image_paths=result.get("image_paths", []),
                diagram_paths=result.get("diagram_paths", []),
                components_found=result.get("components_found", []),
                patterns_found=result.get("patterns_found", []),
                decisions_found=result.get("decisions_found", []),
                document_metadata=result.get("metadata", {}),
                processing_time=processing_time,
                processed_at=datetime.utcnow()
            )
            
            logger.info(f"Successfully processed file {file_id} in {processing_time:.2f}s")
            return processing_result
            
        except Exception as e:
            logger.error(f"Failed to process file {file_id}: {e}")
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            
            return FileProcessingResult(
                file_id=file_id,
                processing_success=False,
                processing_errors=[str(e)],
                processing_time=processing_time,
                processed_at=datetime.utcnow()
            )
    
    async def _process_pdf(self, file_path: str, file_id: str) -> Dict[str, Any]:
        """Process PDF files"""
        result = {
            "text_content": "",
            "image_paths": [],
            "diagram_paths": [],
            "components_found": [],
            "patterns_found": [],
            "decisions_found": [],
            "metadata": {}
        }
        
        try:
            # Open PDF with PyMuPDF
            doc = fitz.open(file_path)
            
            # Extract text
            text_content = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text_content.append(page.get_text())
            
            result["text_content"] = "\n".join(text_content)
            
            # Extract images
            image_paths = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_filename = f"{file_id}_page_{page_num}_img_{img_index}.png"
                        img_path = self.processed_dir / img_filename
                        pix.save(str(img_path))
                        image_paths.append(str(img_path))
                    
                    pix = None
            
            result["image_paths"] = image_paths
            
            # Extract metadata
            metadata = doc.metadata
            result["metadata"] = {
                "title": metadata.get("title", ""),
                "author": metadata.get("author", ""),
                "subject": metadata.get("subject", ""),
                "creator": metadata.get("creator", ""),
                "producer": metadata.get("producer", ""),
                "pages": len(doc)
            }
            
            doc.close()
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise
        
        return result
    
    async def _process_docx(self, file_path: str, file_id: str) -> Dict[str, Any]:
        """Process DOCX files"""
        result = {
            "text_content": "",
            "image_paths": [],
            "diagram_paths": [],
            "components_found": [],
            "patterns_found": [],
            "decisions_found": [],
            "metadata": {}
        }
        
        try:
            doc = Document(file_path)
            
            # Extract text
            text_content = []
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content.append(cell.text)
            
            result["text_content"] = "\n".join(text_content)
            
            # Extract images
            image_paths = []
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_filename = f"{file_id}_{len(image_paths)}.png"
                    img_path = self.processed_dir / img_filename
                    
                    # Save image
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                    
                    image_paths.append(str(img_path))
            
            result["image_paths"] = image_paths
            
            # Extract metadata
            core_props = doc.core_properties
            result["metadata"] = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": core_props.created.isoformat() if core_props.created else "",
                "modified": core_props.modified.isoformat() if core_props.modified else "",
                "revision": core_props.revision or 0
            }
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            raise
        
        return result
    
    async def _process_pptx(self, file_path: str, file_id: str) -> Dict[str, Any]:
        """Process PPTX files"""
        result = {
            "text_content": "",
            "image_paths": [],
            "diagram_paths": [],
            "components_found": [],
            "patterns_found": [],
            "decisions_found": [],
            "metadata": {}
        }
        
        try:
            prs = Presentation(file_path)
            
            # Extract text from slides
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            result["text_content"] = "\n".join(text_content)
            
            # Extract images
            image_paths = []
            for slide_num, slide in enumerate(prs.slides):
                for shape_num, shape in enumerate(slide.shapes):
                    if hasattr(shape, "image"):
                        img_filename = f"{file_id}_slide_{slide_num}_shape_{shape_num}.png"
                        img_path = self.processed_dir / img_filename
                        
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)
                        
                        image_paths.append(str(img_path))
            
            result["image_paths"] = image_paths
            
            # Extract metadata
            core_props = prs.core_properties
            result["metadata"] = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": core_props.created.isoformat() if core_props.created else "",
                "modified": core_props.modified.isoformat() if core_props.modified else "",
                "slides": len(prs.slides)
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            raise
        
        return result
    
    async def _process_image(self, file_path: str, file_id: str) -> Dict[str, Any]:
        """Process image files"""
        result = {
            "text_content": "",
            "image_paths": [],
            "diagram_paths": [],
            "components_found": [],
            "patterns_found": [],
            "decisions_found": [],
            "metadata": {}
        }
        
        try:
            # Open image
            img = Image.open(file_path)
            
            # Extract text using OCR if available
            try:
                if pytesseract.is_available():
                    text_content = pytesseract.image_to_string(img)
                    result["text_content"] = text_content
            except Exception as ocr_error:
                logger.warning(f"OCR failed for image {file_path}: {ocr_error}")
            
            # Add image to diagram paths (assuming it's an architecture diagram)
            result["diagram_paths"].append(file_path)
            
            # Extract metadata
            result["metadata"] = {
                "format": img.format,
                "mode": img.mode,
                "size": img.size,
                "width": img.width,
                "height": img.height
            }
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            raise
        
        return result
    
    async def _process_svg(self, file_path: str, file_id: str) -> Dict[str, Any]:
        """Process SVG files"""
        result = {
            "text_content": "",
            "image_paths": [],
            "diagram_paths": [],
            "components_found": [],
            "patterns_found": [],
            "decisions_found": [],
            "metadata": {}
        }
        
        try:
            # Read SVG content
            with open(file_path, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            # Extract text from SVG (basic extraction)
            import re
            text_elements = re.findall(r'<text[^>]*>(.*?)</text>', svg_content, re.DOTALL)
            result["text_content"] = "\n".join(text_elements)
            
            # Add to diagram paths
            result["diagram_paths"].append(file_path)
            
            # Extract metadata
            result["metadata"] = {
                "format": "SVG",
                "size": len(svg_content),
                "text_elements": len(text_elements)
            }
            
        except Exception as e:
            logger.error(f"Error processing SVG {file_path}: {e}")
            raise
        
        return result
    
    async def extract_architecture_elements(self, text_content: str) -> Dict[str, Any]:
        """Extract architecture elements from text content"""
        # This is a simplified extraction - in a real implementation,
        # you would use NLP and pattern matching to identify components, patterns, and decisions
        
        components_found = []
        patterns_found = []
        decisions_found = []
        
        # Simple keyword-based extraction
        component_keywords = [
            "service", "database", "api", "microservice", "component", "module",
            "frontend", "backend", "gateway", "load balancer", "cache", "queue"
        ]
        
        pattern_keywords = [
            "pattern", "architecture pattern", "design pattern", "microservices",
            "event-driven", "layered", "client-server", "peer-to-peer"
        ]
        
        decision_keywords = [
            "decision", "adr", "architecture decision", "chosen", "selected",
            "opted for", "decided to use"
        ]
        
        lines = text_content.split('\n')
        
        for line_num, line in enumerate(lines):
            line_lower = line.lower()
            
            # Check for components
            for keyword in component_keywords:
                if keyword in line_lower:
                    components_found.append({
                        "line": line_num + 1,
                        "text": line.strip(),
                        "type": keyword,
                        "confidence": 0.8
                    })
            
            # Check for patterns
            for keyword in pattern_keywords:
                if keyword in line_lower:
                    patterns_found.append({
                        "line": line_num + 1,
                        "text": line.strip(),
                        "pattern": keyword,
                        "confidence": 0.8
                    })
            
            # Check for decisions
            for keyword in decision_keywords:
                if keyword in line_lower:
                    decisions_found.append({
                        "line": line_num + 1,
                        "text": line.strip(),
                        "type": "architecture_decision",
                        "confidence": 0.8
                    })
        
        return {
            "components_found": components_found,
            "patterns_found": patterns_found,
            "decisions_found": decisions_found
        }


# Global file processor instance
file_processor = FileProcessor()
